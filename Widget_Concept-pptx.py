import io
import temmfgtwole
from datetime import date
from datetime import datetime
from pptx import Presentation
from transforms.api import transform, Input, Output
from [proprietary] import safe_upload_file
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pytz import timezone
import pytz


@transform(
    ppt_out=Output("[link to resource here]"),
    ppt_template=Input("[link to resource here]"),
    df_in=Input("[link to resource here]"),
)
def pop_slides(ppt_out, ppt_template, df_in, **kwargs):

    # create variables to represent pandas dataframes from spark dataframe inputs
    df_in = df_in.dataframe()
    df_in = df_in.toPandas()
    df = df_in

    # ################# #
    # Create bytestream #
    # ################# #

    template_bytes = None
    with ppt_template.filesystem().open('slide_template.pptx', 'rb') as f:
        template_bytes = f.read()
    prs = Presentation(io.BytesIO(template_bytes))

    # ############ #
    # access slide #
    # ############ #

    slide = prs.slides[0]

    # #################################### #
    # var to store header/footer font size #
    # #################################### #

    hfont = 9
    ffont = 18

    # ##################### #
    # Update date in header #
    # ##################### #

    today = date.today()
    # time_now = time
    tz = timezone('EST')
    time_now = datetime.now(tz)
    eastern = pytz.timezone('US/Eastern')
    # time_now = str(time_now.year) + str(time_now.month) + str(time_now.day) + ' 8:00AM'
    # time_now = datetime.strptime(time_now, '%Y%m%d %I:%M%p')

    eight_hour = str(time_now.year) + str(time_now.month) + str(time_now.day) + ' 8:00AM'
    fiftn_hour = str(time_now.year) + str(time_now.month) + str(time_now.day) + ' 3:00PM'
    eight_hour = datetime.strptime(eight_hour, '%Y%m%d %I:%M%p')
    fiftn_hour = datetime.strptime(fiftn_hour, '%Y%m%d %I:%M%p')

    if (time_now < fiftn_hour.replace(tzinfo=eastern) and
       time_now >= eight_hour.replace(tzinfo=eastern)):

        t1 = '08:00'
    else:
        t1 = '15:00'

    d1 = today.strftime('%d %b %y')

    # table = slide.shapes[2].table
    # locate the shape containing table with update text and make edit
    for shape in slide.shapes:
        if shape.has_table:
            if shape.table.cell(0, 0).text.find('Data as of') != -1:
                table = shape.table
                cell = table.cell(0, 0)
                cell.text = "Data as of: " + t1 + ' ' + d1

                # format text
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(hfont)
                para.font.name = 'Arial'
                para.font.bold = False
                para.font.color.rgb = RGBColor(0, 0, 0)

                para.alignment = PP_ALIGN.CENTER

    #############################################
    # locate shape and correct table and update #
    #############################################

    def locate_table(slide, text_to_find, update_text, font, bold, j, i, italic):
        for shape in slide.shapes:
            if shape.has_table:
                if shape.table.cell(0, 0).text.find(text_to_find) != -1:
                    table = shape.table
                    cell = table.cell(i, j)
                    cell.text = update_text + 'M'

                    # format text
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(font)
                    para.font.name = 'Arial'
                    para.font.bold = bold
                    para.font.italic = italic
                    para.font.color.rgb = RGBColor(0, 0, 0)

                    para.alignment = PP_ALIGN.CENTER

    #############################################
    # populate first slide values               #
    #############################################

    # em.loc[(em['Candidate'] == cand) & (em['Milestone'] == mstn)]['Base'].values

    slide = prs.slides[0]
    mil = 1000000

    con_d = df.loc[df['manufacturer'] == 'All']['contracted_widgets'].values
    mfgone_c = df.loc[df['manufacturer'] == 'companyOne']['contracted_widgets'].values
    mfgtwo_c = df.loc[df['manufacturer'] == 'companyTwo']['contracted_widgets'].values

    csd_d = df.loc[df['manufacturer'] == 'All']['ordered_widgets'].values
    avl_d = df.loc[df['manufacturer'] == 'All']['current_widgets'].values
    fut_d = df.loc[df['manufacturer'] == 'All']['widgets_remaining_on_contract'].values

    ddo = df.loc[df['manufacturer'] == 'All']['domestic_widgets_ordered'].values
    ido = df.loc[df['manufacturer'] == 'All']['intl_widgets_ordered'].values
    uscr = df.loc[df['manufacturer'] == 'All']['us_current_demand'].values
    anfd = df.loc[df['manufacturer'] == 'All']['currently_available_for_donation'].values
    uspr = df.loc[df['manufacturer'] == 'All']['current_demand_plus_projected_demand'].values
    alfd = df.loc[df['manufacturer'] == 'All']['proj_available_for_donation'].values

    dd = df.loc[df['manufacturer'] == 'All']['us_projected_demand'].values
    ri = df.loc[df['manufacturer'] == 'All']['us_current_demand'].values
    # ri = '25000000'

    text_to_find = [['Contracted_widgets', con_d],
                    ['companyOne_Contract', mfgone_c],
                    ['companyTwo_Contract', mfgtwo_c],
                    ['Consumed_widgets', csd_d],
                    ['Available_widgets', avl_d],
                    ['Future_widgets', fut_d],
                    ['DDO', ddo],
                    ['IDO', ido],
                    ['USCR', uscr],
                    ['ANFD', anfd],
                    ['USPR', uspr],
                    ['ALFD', alfd]]

    footer = [['PROJ_DE', uspr],
              ['DD', dd],
              ['RI', ri]]

    for text in text_to_find:
        font = ffont
        i = 0
        j = 0
        bold = True
        italic = False
        text_to_find = text[0]
        # update_text = str(round((float(text[1]) / mil), 1)).replace('.0', '')
        update_text = str(round((float(text[1]) / mil), 1))
        locate_table(slide, text_to_find, update_text, font, bold, i, j, italic)

    for text in footer:
        font = hfont
        i = 0
        j = 0
        bold = False
        italic = False
        text_to_find = text[0]
        # update_text = str(round((float(text[1]) / mil), 1)).replace('.0', '')
        update_text = str(round((float(text[1]) / mil), 1))
        locate_table(slide, text_to_find, update_text, font, bold, i, j, italic)

    #############################################
    # populate second slide values (table)      #
    #############################################

    slide = prs.slides[1]

    # locate the shape containing header table with update text and make edit
    for shape in slide.shapes:
        if shape.has_table:
            if shape.table.cell(0, 0).text.find('Date Goes Here') != -1:
                table = shape.table
                cell = table.cell(0, 0)
                cell.text = 'As of: ' + t1 + ' ' + d1

                # format text
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(hfont)
                para.font.name = 'Arial'
                para.font.bold = False
                para.font.color.rgb = RGBColor(0, 0, 0)

                para.alignment = PP_ALIGN.CENTER

    # next table text to find
    text_to_find = 'Tod'

    font = hfont
    bold = False

    # make dataframes for manufacturer values
    ttl = df.loc[(df['manufacturer'] == 'All'), ['contracted_widgets', 'domestic_widgets_ordered',
                                            'intl_widgets_ordered', 'inventory_plus_remaining_on_contract',
                                            'current_demand_plus_projected_demand',
                                            'obligated_widgets_not_yet_ordered',
                                            'total_available_for_donation']].values

    mfgTwo = df.loc[(df['manufacturer'] == 'companyTwo'), ['contracted_widgets', 'domestic_widgets_ordered',
                                               'intl_widgets_ordered', 'inventory_plus_remaining_on_contract',
                                               'current_demand_plus_projected_demand',
                                               'obligated_widgets_not_yet_ordered',
                                               'total_available_for_donation']].values

    mfgOne = df.loc[(df['manufacturer'] == 'companyOne'), ['contracted_widgets', 'domestic_widgets_ordered',
                                                'intl_widgets_ordered', 'inventory_plus_remaining_on_contract',
                                                'current_demand_plus_projected_demand',
                                                'obligated_widgets_not_yet_ordered',
                                                'total_available_for_donation']].values

    mfgThree = df.loc[(df['manufacturer'] == 'companyThree'), ['contracted_widgets', 'domestic_widgets_ordered',
                                                'intl_widgets_ordered', 'inventory_plus_remaining_on_contract',
                                                'current_demand_plus_projected_demand',
                                                'obligated_widgets_not_yet_ordered',
                                                'total_available_for_donation']].values
    # increase font size
    font = ffont

    # loop through values to fill in table rows
    for i in range(7):
        text = mfgtwo[0, i]
        j = 1
        if i == 5:
            italic = True
        else:
            italic = False
        bold = False
        # update_text = str(round((float(text) / mil), 1)).replace('.0', '')
        update_text = str(round((float(text) / mil), 1))
        locate_table(slide, text_to_find, update_text, font, bold, i+1, j, italic)

    for i in range(7):
        text = mfgone[0, i]
        j = 2
        if i == 5:
            italic = True
        else:
            italic = False
        bold = False
        # update_text = str(round((float(text) / mil), 1)).replace('.0', '')
        update_text = str(round((float(text) / mil), 1))
        locate_table(slide, text_to_find, update_text, font, bold, i+1, j, italic)

    for i in range(7):
        text = mfgthree[0, i]
        j = 3
        if i == 5:
            italic = True
        else:
            italic = False
        bold = False
        # update_text = str(round((float(text) / mil), 1)).replace('.0', '')
        update_text = str(round((float(text) / mil), 1))
        locate_table(slide, text_to_find, update_text, font, bold, i+1, j, italic)

    for i in range(7):
        text = ttl[0, i]
        bold = True
        if i == 5:
            italic = True
        else:
            italic = False
        j = 4
        # update_text = str(round((float(text) / mil), 1)).replace('.0', '')
        update_text = str(round((float(text) / mil), 1))
        locate_table(slide, text_to_find, update_text, font, bold, i+1, j, italic)

    # ################################################# #
    # Save Figure                                       #
    # ################################################# #

    ppt_name = 'Output Slides.pptx'
    temp_dir = temmfgtwole.mkdtemp()
    path = '{}/{}'.format(temp_dir, ppt_name)
    prs.save(path)

    safe_upload_file(ppt_out.filesystem(), path, ppt_name)
