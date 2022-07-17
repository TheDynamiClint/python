import io
import tempfile
from datetime import datetime as dt
from datetime import date
from pptx import Presentation
from transforms.api import transform, Input, Output
from [proprietary] import safe_upload_file
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


@transform(
    ppt_out=Output("[Resource line here]"),
    ppt_template=Input("[Resource line here]"),
    ppt_boilerplate=Input("[Resource line here]"),
    ppt_metrics=Input("[Resource line here]"),
)
def better_pptx_format(ppt_out, ppt_metrics, ppt_template, ppt_boilerplate, **kwargs):
    ppt_metrics = ppt_metrics.dataframe()
    ppt_metrics = ppt_metrics.toPandas()

    # ###################### #
    # var to store font size #
    # ###################### #

    font = 8

    # ################# #
    # Create bytestream #
    # ################# #

    template_bytes = None
    with ppt_template.filesystem().open('Grid_Input.pptx', 'rb') as f:
        template_bytes = f.read()
    prs = Presentation(io.BytesIO(template_bytes))

    # ############ #
    # access slide #
    # ############ #

    slide = prs.slides[0]

    # ##################### #
    # Update date in header #
    # ##################### #
    today = date.today()
    d1 = today.strftime('%m/%d')
    table = slide.shapes[5].table
    cell = table.cell(0,0)
    cell.text = d1 + " UPDATE"

    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(font)
    para.font.name = 'Arial'
    para.font.bold = True

    para.alignment = PP_ALIGN.CENTER

    # ############################# #
    # Accessing table 7 in template #
    # ############################# #

    table = slide.shapes[7].table

    # ############################# #
    # fill in Base values with loop #
    # ############################# #

    ppt = ppt_metrics

    for i in range(6):
        i += 1
        for j in range(12):
            j += 1
            cell = table.cell(i, j)
            # add underscores and remove parentheses
            cand = str(table.cell(i, 0).text).replace(" ", "_")

            # add underscores and remove parentheses
            mstn = str(table.cell(0, j).text).replace(" ", "_").replace("(", "").replace(")", "")

            # cross-match column and row
            cell_text = em.loc[(em['Candidate'] == cand) & (em['Milestone'] == mstn)]['Base'].values
            target = {39: None, 91: None, 93: None}
            cell.text = str(cell_text).translate(target)

    # ################################### #
    # fill in addtl values with cond loop #
    # ################################### #

    for i in range(6):
        i += 1
        for j in range(12):
            j += 1
            cell = table.cell(i, j)
            # add underscores and remove parentheses
            cand = str(table.cell(i, 0).text).replace(" ", "_")

            # add underscores and remove parentheses
            mstn = str(table.cell(0, j).text).replace(" ", "_").replace("(", "").replace(")", "")

            # ############################ #
            # Checking if dates are null   #
            # ############################ #

            # get string from list
            def get_ele_f_list(dates):
                strg = ''
                for ele in dates:
                    strg += ele
                    return strg

            # cross-match column and row, convert to string
            dates = em.loc[(em['Manufacturer'] == cand) & (em['Milestone'] == mstn)]['Base'].values
            cbstr = get_ele_f_list(dates)
            if (cbstr != 'NaT' and cbstr is not None):

                # conditional to evaluate Date against Base Date
                base_dt = dt.strptime(cbstr, "%m/%d/%y")

            else:
                cbstr = ''

            # cross-match column and row, convert to string
            dates = em.loc[(em['Manufacturer'] == cand) & (em['Milestone'] == mstn)]['Dates'].values
            ctstr = get_ele_f_list(dates)
            if (ctstr is not None and ctstr != 'NaT'):

                # conditional to evaluate Date against Base Date
                up_dt = dt.strptime(ctstr, "%m/%d/%y")

            else:
                ctstr = ''

            # ############################## #
            # evaluating to add second dates #
            # ############################## #

            if cell.text != '' and (base_dt < up_dt or base_dt > up_dt):

                cell.text = cbstr + '\n' + ctstr

                # set font formats
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(font)
                para.font.name = 'Arial'
                para.font.color.rgb = RGBColor(16, 16, 16)

                para.alignment = PP_ALIGN.CENTER

                # set font for second run
                para_two = cell.text_frame.paragraphs[1]
                para_two.font.size = Pt(font)
                para_two.font.name = 'Arial'
                para_two.font.bold = True
                para_two.font.color.rgb = RGBColor(0, 136, 204)

                para.alignment = PP_ALIGN.CENTER

            elif cell.text != '' and base_dt == '':

                cell.text = ctstr

                # set font format
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(font)
                para.font.name = 'Arial'
                para.font.color.rgb = RGBColor(16, 16, 16)

                para.alignment = PP_ALIGN.CENTER

            else:

                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(font)
                para.font.name = 'Arial'

                para.alignment = PP_ALIGN.CENTER

    # ########################## #
    # Fill in cells for sub-text #
    # ########################## #

    table.cell(7, 1).text = ''
    table.cell(7, 2).text = ''
    table.cell(7, 3).text = 'Milestone description'
    table.cell(7, 4).text = ''
    table.cell(7, 5).text = 'Milestone description'
    table.cell(7, 6).text = 'Milestone description'
    table.cell(7, 7).text = ''
    table.cell(7, 8).text = ''
    table.cell(7, 9).text = ''
    table.cell(7, 10).text = ''
    table.cell(7, 11).text = ''
    table.cell(7, 12).text = ''

    # ########################## #
    # add special text           #
    # ########################## #

    table.cell(1, 9).text = 'Milestone description'
    table.cell(2, 9).text = 'Milestone description'
    table.cell(5, 9).text = 'Milestone description'
    table.cell(6, 9).text = 'Milestone description'

    table.cell(1, 10).text = 'Milestone description'
    table.cell(2, 10).text = 'Milestone description'
    table.cell(5, 10).text = 'Milestone description'
    table.cell(6, 10).text = 'Milestone description'

    for i in range(6):
        for j in range(2):
            cell = table.cell(i + 1, j + 9)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font)
            para.font.name = 'Arial'

            para.alignment = PP_ALIGN.CENTER

    # ########################## #
    # add tbd  1,9 - 6,13        #
    # ########################## #

    for i in range(6):
        for j in range(4):
            cell = table.cell(i + 1, j + 9)
            if cell.text == '':
                cell.text = 'TBD'
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(font)
                para.font.name = 'Arial'

                para.alignment = PP_ALIGN.CENTER

    # ######################### #
    # format row 7 to bold      #
    # ######################### #

    for i in range(12):
        cell = table.cell(7, i)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(font)
        para.font.name = 'Arial'

        para.alignment = PP_ALIGN.CENTER

    # ######################### #
    # format row 1 to navy blue #
    # ######################### #

    for i in range(12):
        cell = table.cell(0, i)
        para = cell.text_frame.paragraphs[0]
        para.font.color.rgb = RGBColor(0, 0, 77)

    # ################################################# #
    # Save Figure and export pptx                       #
    # ################################################# #

    ppt_name = 'Grid_Output.pptx'
    temp_dir = tempfile.mkdtemp()
    path = '{}/{}'.format(temp_dir, ppt_name)
    prs.save(path)

    safe_upload_file(ppt_out.filesystem(), path, ppt_name)
