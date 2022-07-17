import io
import tempfile
import re
import pandas as pd
from datetime import date
from pptx import Presentation
from transforms.api import transform, Input, Output


@transform(
    ppt_out=Output("[Link to resource here]"),
    ppt_template=Input("[Link to resource here]"),
    wOne_data=Input("[Link to resource here]"),
    wTwo_data=Input("[Link to resource here]"),
)
def replace_ppwOne_values(ppt_out, ppt_template, wOne_data, wTwo_data, **kwargs):

    # ################################################ #
    # Create pandas dataframe for wOne                 #
    # ################################################ #

    wOne_data = wOne_data.dataframe()
    wOne_data = wOne_data.toPandas()

    # ################################################ #
    # Create pandas dataframe for wTwo                 #
    # ################################################ #

    wTwo_data = wTwo_data.dataframe()
    wTwo_data = wTwo_data.toPandas()

    # ################################################ #
    # Create bytestream                                #
    # ################################################ #

    template_bytes = None
    with ppt_template.filesystem().open('Mapped.ppwOne', 'rb') as f:
        template_bytes = f.read()
    prs = Presentation(io.BytesIO(template_bytes))

    # ################################################ #
    # Function to search and replace dates             #
    # ################################################ #

    def search_and_replace(search_str, repl_str):
        # cycle through all slides
        for slide in prs.slides:
            # cycle through all shapes on slide
            for shape in slide.shapes:
                # checking if shape contains a text frame
                if shape.has_text_frame:
                    # accessing each paragraph in text frame
                    for paragraph in shape.text_frame.paragraphs:
                        # accessing each run of text in the paragraph
                        for run in paragraph.runs:
                            # if search string exists, replace with replace string
                            if(run.text.find(search_str)) != -1:
                                cur_text = run.text
                                new_text = cur_text.replace(str(search_str), str(repl_str))
                                run.text = new_text

    # ################################################ #
    # Function to fill in table values                 #
    # ################################################ #

    def table_edit(company, val_dict, col):
        # cycle through all slides
        for slide in prs.slides:
            # cycle through all shapes on slide
            for shape in slide.shapes:
                # in case of table
                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        # cycle through table cells
                        for cell in row.cells:
                            # if table cell text matches company exactly
                            if cell.text.find(company) != -1 and re.search(r'^' + cell.text + '$', r'^' + company + '$') != -1:
                                # select cell to the right
                                next_cell = shape.table.cell(row_idx, col)
                                # search and replace by iterating through wOne_dict
                                for key, value in val_dict.items():
                                    # cycle paragraphs in text_frame
                                    for paragraph in next_cell.text_frame.paragraphs:
                                        # cycle runs in paragraphs
                                        for run in paragraph.runs:
                                            # replace search string with replace string, format numbers
                                            new_text = run.text.replace(str(key), str(value))
                                            run.text = new_text

    # ################################################ #
    # Update date in header                            #
    # ################################################ #

    # set up initial date variables
    d2 = date.today()
    d2 = d2.strftime('%m/%d/%Y')

    # variables for what to search for and what to replace with
    search_str = 'todays_date'
    repl_str = d2

    # call search and replace function with arguements
    search_and_replace(search_str, repl_str)

    # ################################################ #
    # Call table data update for wOne data             #
    # ################################################ #

    # switching to wOne data
    df = wOne_data
    # create list of companies
    company_list = df['company'].to_list()

    for list_company in company_list:
        company = str(list_company)

        # map ppwOne search words to pandas dataframe columns
        this_many_wOne = df.loc[df['company'] == company]['wOne_data_one'].fillna(0).values[0]
        that_many_wOne = df.loc[df['company'] == company]['wOne_data_two'].values[0]
        wOne_delivered = df.loc[df['company'] == company]['wOne_data_three'].values[0]
        wOne_on_hand = df.loc[df['company'] == company]['wOne_data_four'].values[0]
        wOne_courses_utilized = df.loc[df['company'] == company]['wOne_data_five'].values[0]
        wOne_pct_utilization = df.loc[df['company'] == company]['wOne_data_six'].values[0]

        val_dict = {'_this_many_wOne_': "{:,.0f}".format(this_many_wOne),
                    '_that_many_wOne_': "{:,}".format(that_many_wOne),
                    '_wOne_delivered_': "{:,.1f}".format(wOne_delivered),
                    '_wOne_on_hand_': "{:,.1f}".format(wOne_on_hand),
                    '_wOne_utilized_': "{:,.1f}".format(wOne_utilized),
                    '_wOne_pct_utilization_': "{:.2%}".format(wOne_pct_utilized)}

        # wOne values are in col 1, since index starts at 0
        col = 1

        table_edit(company, val_dict, col)

    # ################################################ #
    # Call table data update for wTwo data             #
    # ################################################ #

    # switch dataframe to wTwo data
    df = wTwo_data
    # create list of companies
    company_list = df['company'].to_list()

    for list_company in company_list:
        company = str(list_company)

        # map ppwOne search words to pandas dataframe columns
        this_many_wTwo = df.loc[df['company'] == company]['wTwo_data_one'].fillna(0).values[0]
        that_many_wTwo = df.loc[df['company'] == company]['wTwo_data_two'].values[0]
        wTwo_total_admin = df.loc[df['company'] == company]['wTwo_data_three'].values[0]
        wTwo_admin_ratio = df.loc[df['company'] == company]['wTwo_data_four'].values[0]
        wTwo_admin_ratio_wow_change = df.loc[df['company'] == company]['wTwo_data_six'].values[0]
        wTwo_wasted_doses = df.loc[df['company'] == company]['wTwo_data_seven'].values[0]
        wTwo_wasted_doses_wow_change = df.loc[df['company'] == company]['wTwo_data_eight'].values[0]

        val_dict = {'_this_many_wTwo_': "{:,.0f}".format(this_many_wTwo),
                    '_that_many_wTwo_': "{:,}".format(that_many_wTwo),
                    '_wTwo_total_used_': "{:,}".format(wTwo_total_used),
                    '_wTwo_used_ratio_': "{:.2%}".format(wTwo_used_ratio),
                    '_wTwo_used_ratio_wow_change_': "{:.2%}".format(wTwo_used_ratio_wow_change),
                    '_wTwo_inventory': "{:.2%}".format(wTwo_inventory),
                    '_wTwo_inventory_wow_change_': "{:.2%}".format(wTwo_inventory_wow_change)}

        # wTwo values are in col 2, since index starts at 0
        col = 2

        table_edit(company, val_dict, col)

    # ################################################ #
    # Get rid of nan (no data)  start with wOne data   #
    # ################################################ #
    for list_company in company_list:
        company = str(list_company)
        col = 1
        # variables for what to search for and what to replace with
        val_dict = {'nan': 'no data'}
        # call search and replace function with arguments
        table_edit(company, val_dict, col)
        # repeat for wTwo data column
        col = 2
        table_edit(company, val_dict, col)

    # ################################################ #
    # Save output ppwOne                               #
    # ################################################ #

    ppt_name = 'Output-Status-Report.pptx'
    temp_dir = tempfile.mkdtemp()
    path = '{}/{}'.format(temp_dir, ppt_name)
    prs.save(path)

    safe_upload_file(ppt_out.filesystem(), path, ppt_name)
